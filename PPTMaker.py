from pptx import Presentation
import wikipedia
from simple_image_download import simple_image_download as simp
from time import sleep

userInput = ''
topics = []
response = simp.simple_image_download
PPTName = input("Name of PowerPoint: ")

print('Enter topics of slides: ')
while userInput != 'q':
    userInput = input('Topic of slide ')
    if userInput != 'q':
        topics.append([userInput, wikipedia.summary(userInput, sentences=3)])
        response().download(userInput , 1)

def changeChars(word, replaced, replacer):
    newWord = ""
    for char in word:
        if char == replaced:
            newWord += replacer
        else:
            newWord += char
    return newWord

prs = Presentation()
for topic in topics:
    title_slide_layout = prs.slide_layouts[8]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    placeholder = slide.placeholders[1]
    folder = changeChars(topic[0], " ", "_")
    pic = placeholder.insert_picture('simple_images/' + folder + '/' + topic[0] + '_1' + '.jpeg')
    sub = slide.placeholders[2]
    title.text = topic[0].title()
    sub.text = topic[1]
prs.save(PPTName + '.pptx')
